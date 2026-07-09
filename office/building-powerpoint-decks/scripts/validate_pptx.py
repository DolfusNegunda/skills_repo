"""Validate a .pptx deck before it ships.

Deterministic quality gate for a produced deck. Catches the failure modes that make
a deck look unfinished, prints a machine-readable JSON report, and exits non-zero on
ERRORS so the caller can fix and re-run (produce -> validate -> fix -> re-validate).

Checks:
  ERROR   - Leftover placeholder text (lorem/ipsum, TBD, TODO, FIXME, XXX,
            PLACEHOLDER, and {{ }} / {% %} template tags) on any slide or in notes.
  WARNING - Empty slides (no title and no text/table content).
  WARNING - Slides with no title (harder to navigate; check intent).
  WARNING - Very text-dense slides (> ~60 words), a wall-of-text smell.

Usage:
    python scripts/validate_pptx.py path/to/deck.pptx
"""
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PLACEHOLDER = re.compile(
    r"\bTBD\b|\bTODO\b|\bFIXME\b|\bXXX+\b|\bPLACEHOLDER\b|lorem\s*ipsum|{{.*?}}|{%.*?%}",
    re.IGNORECASE,
)
DENSE_WORDS = 60


def slide_texts(slide):
    """(title, [body strings], has_table, notes) for one slide."""
    title = slide.shapes.title.text.strip() if slide.shapes.title else ""
    body, has_table = [], False
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if shape.has_table:
            has_table = True
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        body.append(cell.text.strip())
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        elif shape.has_text_frame and shape.text_frame.text.strip():
            body.append(shape.text_frame.text.strip())
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    return title, body, has_table, notes


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: python scripts/validate_pptx.py <deck.pptx>")
    path = Path(sys.argv[1])
    if not path.exists():
        sys.exit(f"file not found: {path}")

    prs = Presentation(str(path))
    errors, warnings = [], []
    placeholders = []

    for idx, slide in enumerate(prs.slides, start=1):
        title, body, has_table, notes = slide_texts(slide)
        for chunk in [title, *body, notes]:
            m = PLACEHOLDER.search(chunk or "")
            if m:
                placeholders.append(f"slide {idx}: {m.group(0)!r}")
        if not title and not body and not has_table:
            warnings.append(f"slide {idx}: empty (no title or content)")
        elif not title:
            warnings.append(f"slide {idx}: no title")
        words = sum(len(b.split()) for b in body)
        if words > DENSE_WORDS:
            warnings.append(f"slide {idx}: dense (~{words} words) — wall-of-text risk")

    if placeholders:
        errors.append(f"Leftover placeholder text: {placeholders}")

    report = {
        "file": str(path),
        "status": "OK" if not errors else "FAILED",
        "errors": errors,
        "warnings": warnings,
        "n_slides": len(prs.slides._sldIdLst),
    }
    print(json.dumps(report, indent=2, ensure_ascii=False))
    sys.exit(0 if not errors else 1)


if __name__ == "__main__":
    main()
