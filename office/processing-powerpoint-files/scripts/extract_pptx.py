"""Extract a .pptx into clean, per-slide, model-readable JSON.

Deterministic ingestion for PowerPoint: for each slide it captures the title, the
body text in top-to-bottom / left-to-right reading order, tables, and the speaker
notes (the part a naive dump misses), plus a fidelity block (slide count, images,
notes present). Prints JSON; exits 0 on a successful read.

This is an EXTRACTOR (ingest + fidelity self-check), not a validate->fix loop.

Usage:
    python scripts/extract_pptx.py path/to/deck.pptx
"""
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def shape_key(shape):
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    return (top, left)


def table_rows(shape):
    tbl = shape.table
    return [[cell.text.strip() for cell in row.cells] for row in tbl.rows]


def main():
    if len(sys.argv) != 2:
        sys.exit("usage: python scripts/extract_pptx.py <deck.pptx>")
    path = Path(sys.argv[1])
    if not path.exists():
        sys.exit(f"file not found: {path}")

    prs = Presentation(str(path))
    slides, n_images = [], 0
    slides_with_notes = 0

    for idx, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else None
        body, tables = [], []
        for shape in sorted(slide.shapes, key=shape_key):
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n_images += 1
                continue
            if shape.has_table:
                tables.append(table_rows(shape))
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    body.append(text)
        notes = None
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip() or None
            if notes:
                slides_with_notes += 1
        slides.append({
            "slide": idx,
            "title": title,
            "body": body,
            "tables": tables,
            "notes": notes,
        })

    result = {
        "file": str(path),
        "format": "pptx",
        "slides": slides,
        "fidelity": {
            "n_slides": len(slides),
            "n_images": n_images,
            "slides_with_notes": slides_with_notes,
            "warnings": (
                ["No slides found."] if not slides else []
            ),
        },
    }
    print(json.dumps(result, indent=2, ensure_ascii=False))
    sys.exit(0)


if __name__ == "__main__":
    main()
