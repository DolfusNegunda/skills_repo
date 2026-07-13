"""Shared helpers for the document-template engine.

The engine turns a client's existing file into a reusable *template* (same layout,
fonts, logos, styles — with the variable text swapped for `{{ placeholder }}` tags)
plus a machine-readable *manifest* describing every placeholder. A later fill step
reads the manifest, drops in content, and re-emits a finished document.

Design choices that live here:

* **Placeholder convention** — one simple form everywhere: ``{{ field_name }}``.
  No loop/branch syntax in the template. List and table expansion is driven by the
  *manifest field type* at fill time, not by tags in the document. This keeps
  templatize (which only has ONE example to learn from) from having to synthesize
  control-flow, and makes DOCX and PPTX behave identically.

* **Run-aware, text-based replacement** — we replace the *text* of a value wherever
  it appears, preserving the formatting of the run it lives in. When a value spans
  several runs we consolidate them into the first run (keeping its formatting) so
  the tag always lands in a single run — the one failure mode that otherwise breaks
  in-place fills. Replacing every occurrence is usually what you want (a client name
  appears in many places and should all become ``{{ client_name }}``).

This module is format-agnostic where it can be: it exposes paragraph iterators for
DOCX and PPTX that both yield objects with a ``.runs`` list of runs that each carry
a settable ``.text`` — so the scan/replace/render logic is written once.
"""
from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, asdict, field as dc_field
from pathlib import Path

# ── Paths ──────────────────────────────────────────────────────────────────────
# scripts/ -> skill root -> registry/ is the version-controlled template gallery.
# Override with $TEMPLATE_REGISTRY to point at a shared gallery outside the repo
# (or to isolate tests).
SKILL_ROOT = Path(__file__).resolve().parents[1]
REGISTRY = Path(os.environ.get("TEMPLATE_REGISTRY", SKILL_ROOT / "registry"))

# ── Placeholder convention ───────────────────────────────────────────────────--
# {{ field_name }} — optional surrounding whitespace; names are [a-z0-9_].
PLACEHOLDER_RE = re.compile(r"\{\{\s*([a-zA-Z0-9_]+)\s*\}\}")
# Any leftover tag (used by the validator to detect an unfilled/partial fill).
ANY_TAG_RE = re.compile(r"\{\{.*?\}\}|\{%.*?%\}")

SUPPORTED_FORMATS = ("docx", "pptx")


def detect_format(path: str | Path) -> str:
    """Map a file path to a supported engine format by extension."""
    ext = Path(path).suffix.lower().lstrip(".")
    if ext not in SUPPORTED_FORMATS:
        raise ValueError(
            f"Unsupported format '.{ext}'. This engine templatizes {SUPPORTED_FORMATS}. "
            "PDF is an export target (fill a DOCX/PPTX then export), not a templatize source."
        )
    return ext


def slugify(text: str) -> str:
    """A filesystem/identifier-safe kebab/snake slug (lowercase, _-separated words)."""
    text = re.sub(r"[^\w\s-]", "", str(text).strip().lower())
    text = re.sub(r"[\s-]+", "_", text)
    return text.strip("_") or "unnamed"


def placeholder(name: str) -> str:
    """Render the canonical tag string for a field name."""
    return "{{ " + name + " }}"


# ── Manifest model ───────────────────────────────────────────────────────────--
@dataclass
class Field:
    """One variable slot in a template."""
    name: str                       # snake_case identifier, e.g. "client_name"
    type: str = "text"              # text | list | table_rows
    example: str = ""               # the value seen in the source (a fill sample)
    guidance: str = ""              # how to fill it, for humans/agents
    required: bool = True

    def to_dict(self) -> dict:
        return asdict(self)


@dataclass
class Manifest:
    """Everything a future fill needs, without re-reading the whole document."""
    template_id: str                # "<client>/<doc_type>"
    client: str
    doc_type: str
    format: str                     # docx | pptx
    template_file: str              # relative name inside the template dir
    source_file: str = ""           # original client file this was learned from
    version: str = "1.0.0"
    owner: str = ""
    created: str = ""               # ISO date; caller stamps it (no clock in scripts)
    changelog: list = dc_field(default_factory=list)
    fields: list = dc_field(default_factory=list)  # list[Field]

    def to_dict(self) -> dict:
        d = asdict(self)
        d["fields"] = [f.to_dict() if isinstance(f, Field) else f for f in self.fields]
        return d

    def field_names(self) -> list[str]:
        return [f["name"] if isinstance(f, dict) else f.name for f in self.fields]


def save_manifest(manifest: Manifest, path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(manifest.to_dict(), indent=2, ensure_ascii=False),
                    encoding="utf-8")


def load_manifest(path: Path) -> dict:
    return json.loads(Path(path).read_text(encoding="utf-8"))


# ── Registry (the template gallery) ─────────────────────────────────────────--
def template_dir(client: str, doc_type: str) -> Path:
    """registry/<client>/<doc_type>/ — where a template + its manifest live."""
    return REGISTRY / slugify(client) / slugify(doc_type)


def find_template(client: str, doc_type: str) -> tuple[Path, dict]:
    """Resolve a registered template. Returns (template_path, manifest_dict)."""
    d = template_dir(client, doc_type)
    man_path = d / "manifest.json"
    if not man_path.exists():
        raise FileNotFoundError(
            f"No template registered for client='{client}', doc_type='{doc_type}'. "
            f"Expected {man_path}. Run registry.py list to see what exists."
        )
    manifest = load_manifest(man_path)
    tpl_path = d / manifest["template_file"]
    if not tpl_path.exists():
        raise FileNotFoundError(f"Manifest references missing template file: {tpl_path}")
    return tpl_path, manifest


# ── Format-agnostic paragraph iteration ─────────────────────────────────────--
# Both python-docx and python-pptx paragraphs expose `.runs` where each run has a
# settable `.text`. We yield paragraph objects so scan/replace/render is written once.

def _iter_docx_elem(element, parent):
    """Yield paragraphs under an lxml element (body or table cell), in document
    order, recursing into nested tables."""
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    for child in element.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, parent)
        elif child.tag == qn("w:tbl"):
            table = Table(child, parent)
            for row in table.rows:
                for cell in row.cells:
                    yield from _iter_docx_elem(cell._element, cell)


def iter_docx_paragraphs(doc):
    """Every paragraph in a .docx: body, tables (recursive), and headers/footers."""
    yield from _iter_docx_elem(doc.element.body, doc._body)
    for section in doc.sections:
        for hf in (section.header, section.footer,
                   section.first_page_header, section.first_page_footer,
                   section.even_page_header, section.even_page_footer):
            for para in hf.paragraphs:
                yield para
            for table in hf.tables:
                for row in table.rows:
                    for cell in row.cells:
                        yield from _iter_docx_elem(cell._element, cell)


def iter_pptx_paragraphs(prs):
    """Every text paragraph in a .pptx across slides, shapes (incl. groups) & tables."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def walk_shapes(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk_shapes(shape.shapes)
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield para
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield para

    for slide in prs.slides:
        yield from walk_shapes(slide.shapes)
        if slide.has_notes_slide:
            for para in slide.notes_slide.notes_text_frame.paragraphs:
                yield para


def para_text(paragraph) -> str:
    """Concatenated text of a paragraph's runs (works for docx and pptx)."""
    return "".join(run.text for run in paragraph.runs)


def replace_in_paragraph(paragraph, old: str, new: str) -> int:
    """Replace every occurrence of `old` with `new` inside one paragraph.

    Preserves run formatting. If `old` sits within a single run we edit that run in
    place. If it spans runs we consolidate the paragraph's runs into the first
    (keeping the first run's formatting) so `new` — critically, a ``{{ tag }}`` —
    lands in exactly one run. Returns the number of replacements made.
    """
    if not old:
        return 0
    runs = paragraph.runs
    if not runs:
        return 0

    # Fast path: value contained within a single run.
    n = 0
    for run in runs:
        if old in run.text:
            count = run.text.count(old)
            run.text = run.text.replace(old, new)
            n += count
    if n:
        return n

    # Slow path: value spans runs -> consolidate into the first run, then replace.
    full = "".join(r.text for r in runs)
    if old not in full:
        return 0
    count = full.count(old)
    runs[0].text = full.replace(old, new)
    for r in runs[1:]:
        r.text = ""
    return count
