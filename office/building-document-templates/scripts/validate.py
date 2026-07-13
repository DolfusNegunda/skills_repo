"""Gate a filled document before it ships.

Fails loudly (non-zero exit) with specific messages so problems are fixable:
  * No leftover ``{{ tag }}`` / ``{% ... %}`` anywhere (body, tables, headers,
    footers, slides, notes) — the classic "shipped with a placeholder in it" bug.
  * The output still has structure — non-empty, and (vs. the template, if given)
    the same slide count / no lost sections.

Prints a JSON report; exit 0 only when status == "OK".

Usage:
    python scripts/validate.py out/acme-q4.docx
    python scripts/validate.py out/deck.pptx --template registry/acme/board-deck/template.pptx
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import common as C


def docx_report(path: Path, template: Path | None):
    from docx import Document
    doc = Document(str(path))
    texts = [C.para_text(p) for p in C.iter_docx_paragraphs(doc)]
    errors, checks = [], {}

    leftover = sorted({m.group(0) for t in texts for m in C.ANY_TAG_RE.finditer(t)})
    if leftover:
        errors.append(f"Unfilled template tags remain: {leftover}")
    checks["nonempty_paragraphs"] = sum(1 for t in texts if t.strip())
    if checks["nonempty_paragraphs"] == 0:
        errors.append("Document has no text content.")

    if template and template.exists():
        tdoc = Document(str(template))
        t_sections = len(tdoc.sections)
        checks["sections"] = len(doc.sections)
        if len(doc.sections) != t_sections:
            errors.append(f"Section count changed: template {t_sections} -> output {len(doc.sections)}")
    elif template:
        errors.append(f"--template not found: {template} (structure check could not run)")
    return errors, checks


def pptx_report(path: Path, template: Path | None):
    from pptx import Presentation
    prs = Presentation(str(path))
    texts = [C.para_text(p) for p in C.iter_pptx_paragraphs(prs)]
    errors, checks = [], {}

    leftover = sorted({m.group(0) for t in texts for m in C.ANY_TAG_RE.finditer(t)})
    if leftover:
        errors.append(f"Unfilled template tags remain: {leftover}")
    checks["n_slides"] = len(prs.slides)
    if checks["n_slides"] == 0:
        errors.append("Deck has no slides.")

    if template and template.exists():
        tprs = Presentation(str(template))
        t_slides = len(tprs.slides)
        if checks["n_slides"] != t_slides:
            errors.append(f"Slide count changed: template {t_slides} -> output {checks['n_slides']}")
    elif template:
        errors.append(f"--template not found: {template} (structure check could not run)")
    return errors, checks


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("file")
    ap.add_argument("--template", default=None, help="Compare structure against this template")
    args = ap.parse_args()

    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    template = Path(args.template) if args.template else None
    fmt = C.detect_format(path)

    errors, checks = (docx_report if fmt == "docx" else pptx_report)(path, template)
    report = {
        "file": str(path),
        "format": fmt,
        "status": "OK" if not errors else "FAIL",
        "checks": checks,
        "errors": errors,
    }
    print(json.dumps(report, indent=2, ensure_ascii=False))
    sys.exit(0 if not errors else 1)


if __name__ == "__main__":
    main()
