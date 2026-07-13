"""Turn a client's existing .docx/.pptx into a reusable template + manifest.

Templatizing preserves the file — layout, fonts, logos, styles, slide masters all
stay exactly as the client made them — and only swaps the *variable* text for
``{{ placeholder }}`` tags. It is a TWO-STEP, ASSISTED process because a single
example can't tell you which text is boilerplate and which changes each time
("Acme Q3 2026" — is Acme the client? is Q3 the quarter? both?):

  1. propose : read the file, extract candidate variable values with context and a
               heuristic guess, and write a proposal JSON for a human/agent to edit.
  2. build   : read the (edited) proposal, inject placeholders for everything marked
               keep="variable", and register the template + manifest in the gallery.

Usage:
    python scripts/templatize.py propose --file client.docx --out proposal.json
    # ...review proposal.json: set keep, rename fields, set type (text|list)...
    python scripts/templatize.py build --file client.docx --fields proposal.json \
        --client acme --doc-type quarterly-review \
        --owner you@co.com --created 2026-07-13
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

import common as C

DATE_RE = re.compile(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b|\b\d{4}-\d{2}-\d{2}\b"
                     r"|\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\.?\s+\d{1,2},?\s+\d{4}\b",
                     re.I)
MONEY_RE = re.compile(r"[$£€]\s?\d")
LABEL_RE = re.compile(r"^\s*([A-Za-z][A-Za-z /&-]{1,30}?)\s*[:\-–]\s*(.+\S)\s*$")
QUARTER_RE = re.compile(r"\bQ[1-4]\b", re.I)

# Paragraphs longer than this are treated as prose/boilerplate by default (keep=fixed).
LONG_TEXT = 140


def load(path: Path):
    fmt = C.detect_format(path)
    if fmt == "docx":
        from docx import Document
        doc = Document(str(path))
        return fmt, doc, list(C.iter_docx_paragraphs(doc))
    from pptx import Presentation
    prs = Presentation(str(path))
    return fmt, prs, list(C.iter_pptx_paragraphs(prs))


def suggest_name(value: str, label: str | None, seen: set[str]) -> str:
    """Best-effort snake_case field name from a label, else the value's shape."""
    if label:
        base = C.slugify(label)
    elif DATE_RE.search(value):
        base = "date"
    elif MONEY_RE.search(value):
        base = "amount"
    elif QUARTER_RE.search(value) and len(value) <= 12:
        base = "period"
    elif len(value) <= 40:
        base = C.slugify(value)[:30] or "field"
    else:
        base = "field"
    name, i = base, 2
    while name in seen:
        name, i = f"{base}_{i}", i + 1
    seen.add(name)
    return name


def guess_keep_and_type(value: str) -> tuple[str, str]:
    """Heuristic default: is this variable content, and is it text or a list?"""
    if len(value) > LONG_TEXT:
        return "fixed", "text"       # prose/boilerplate — flip to variable if needed
    return "variable", "text"


def propose(args):
    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    fmt, _doc, paragraphs = load(path)

    # Dedupe by exact value (replace-all semantics): the same client name in 5
    # places becomes ONE field that fills everywhere.
    by_value: dict[str, dict] = {}
    seen_names: set[str] = set()
    for para in paragraphs:
        text = C.para_text(para).strip()
        if not text:
            continue
        m = LABEL_RE.match(text)
        # For "Label: value" lines, the *value* is the candidate; keep the label fixed.
        value, label = (m.group(2).strip(), m.group(1).strip()) if m else (text, None)
        if not value or C.PLACEHOLDER_RE.search(value):
            continue
        entry = by_value.get(value)
        if entry:
            entry["occurrences"] += 1
            if label and label not in entry["labels"]:
                entry["labels"].append(label)
            continue
        keep, ftype = guess_keep_and_type(value)
        by_value[value] = {
            "current_text": value,
            "occurrences": 1,
            "labels": [label] if label else [],
            "context": text if text != value else "",
            "suggest_name": suggest_name(value, label, seen_names),
            "suggest_type": ftype,
            "keep": keep,
        }

    candidates = sorted(by_value.values(), key=lambda c: (-len(c["current_text"])))
    proposal = {
        "format": fmt,
        "source_file": str(path),
        "instructions": (
            "Review each candidate and set keep to one of: 'variable' (filled each "
            "time), 'fixed' (boilerplate that stays), or 'remove' (delete this line — "
            "use for surplus example bullets/rows a single 'list' field will "
            "regenerate). For a repeating list, mark ONE representative line "
            "keep='variable' + suggest_type='list', and mark the other example lines "
            "keep='remove'. Rename suggest_name to a clean snake_case field. Delete "
            "candidates you don't care about. Then: templatize.py build --file <file> "
            "--fields <this>."
        ),
        "candidates": candidates,
    }
    out = Path(args.out)
    out.write_text(json.dumps(proposal, indent=2, ensure_ascii=False), encoding="utf-8")
    n_var = sum(1 for c in candidates if c["keep"] == "variable")
    print(f"Proposed {len(candidates)} candidates ({n_var} variable) -> {out}")
    print("Review/edit it, then run `templatize.py build`.")


def build(args):
    path = Path(args.file)
    if not path.exists():
        sys.exit(f"file not found: {path}")
    proposal = json.loads(Path(args.fields).read_text(encoding="utf-8"))
    variables = [c for c in proposal["candidates"] if c.get("keep") == "variable"]
    removals = [c for c in proposal["candidates"] if c.get("keep") == "remove"]
    if not variables:
        sys.exit("No candidates marked keep='variable' — nothing to templatize.")

    fmt, doc, paragraphs = load(path)
    if fmt != proposal.get("format"):
        sys.exit(f"Proposal format '{proposal.get('format')}' != file format '{fmt}'.")

    # Remove extra example paragraphs first (e.g. surplus bullets/rows a `list` field
    # will regenerate). Match on exact stripped text so we only drop intended lines.
    for c in removals:
        target = c["current_text"].strip()
        for para in list(paragraphs):
            if C.para_text(para).strip() == target:
                para._p.getparent().remove(para._p)
    paragraphs = [p for p in paragraphs if p._p.getparent() is not None]

    # Longest first so a longer value ("Q3 2026") is tagged before a shorter substring
    # ("Q3") could clobber part of it.
    variables.sort(key=lambda c: -len(c["current_text"]))

    fields, warnings = [], []
    for c in variables:
        name = C.slugify(c.get("suggest_name") or c["current_text"])
        tag = C.placeholder(name)
        replaced = 0
        for para in paragraphs:
            replaced += C.replace_in_paragraph(para, c["current_text"], tag)
        if replaced == 0:
            warnings.append(f"'{c['current_text'][:40]}' -> {{{{ {name} }}}}: 0 matches "
                            "(text may span formatting boundaries oddly; check manually)")
        fields.append(C.Field(
            name=name,
            type=c.get("suggest_type", "text"),
            example=c["current_text"],
            guidance=(f"Fills: {', '.join(c.get('labels', []))}" if c.get("labels") else ""),
            required=True,
        ))

    tdir = C.template_dir(args.client, args.doc_type)
    tdir.mkdir(parents=True, exist_ok=True)
    template_file = f"template.{fmt}"
    doc.save(str(tdir / template_file))

    manifest = C.Manifest(
        template_id=f"{C.slugify(args.client)}/{C.slugify(args.doc_type)}",
        client=C.slugify(args.client),
        doc_type=C.slugify(args.doc_type),
        format=fmt,
        template_file=template_file,
        source_file=Path(path).name,
        owner=args.owner,
        created=args.created,
        changelog=[f"{args.created} v1.0.0 templatized from {Path(path).name}"] if args.created else [],
        fields=fields,
    )
    C.save_manifest(manifest, tdir / "manifest.json")

    print(f"Template registered: {manifest.template_id}  ({len(fields)} fields, {fmt})")
    print(f"  {tdir / template_file}")
    print(f"  {tdir / 'manifest.json'}")
    for w in warnings:
        print(f"  WARNING: {w}")


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    sub = ap.add_subparsers(dest="cmd", required=True)

    p = sub.add_parser("propose", help="Extract candidate variable fields for review")
    p.add_argument("--file", required=True, help="Client .docx/.pptx to learn from")
    p.add_argument("--out", default="proposal.json")
    p.set_defaults(func=propose)

    b = sub.add_parser("build", help="Inject placeholders and register the template")
    b.add_argument("--file", required=True, help="Same client file as `propose`")
    b.add_argument("--fields", required=True, help="Reviewed proposal JSON")
    b.add_argument("--client", required=True)
    b.add_argument("--doc-type", required=True)
    b.add_argument("--owner", default="")
    b.add_argument("--created", default="", help="ISO date (scripts have no clock)")
    b.set_defaults(func=build)

    args = ap.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
