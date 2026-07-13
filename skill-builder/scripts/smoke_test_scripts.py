"""Smoke-test every bundled skill script against generated fixtures.

Proves the document scripts actually run and behave as documented (right exit codes,
right verdicts) — reproducibly, in CI or by hand. Generates real .xlsx/.docx/.pptx/
.pdf/.csv fixtures, runs each validator/extractor, and asserts the outcome.

Run from the repo root:
    python skill-builder/scripts/smoke_test_scripts.py

Requires: openpyxl, python-docx, python-pptx, pypdf, reportlab.
Exits non-zero if any script is missing, errors unexpectedly, or gives a wrong verdict.
"""
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
results = []


def run(script_rel, *args, env=None):
    """Run a bundled script; return (exit_code, parsed_json_or_None)."""
    script = ROOT / script_rel
    if not script.exists():
        return None, None
    run_env = None
    if env:
        run_env = os.environ.copy()
        run_env.update({k: str(v) for k, v in env.items()})
    proc = subprocess.run([sys.executable, str(script), *[str(a) for a in args]],
                          capture_output=True, text=True, env=run_env)
    data = None
    try:
        data = json.loads(proc.stdout)
    except (ValueError, json.JSONDecodeError):
        pass
    return proc.returncode, data


def check(name, cond):
    results.append((name, bool(cond)))
    print(f"  {'PASS' if cond else 'FAIL'}  {name}")


def make_fixtures(tmp):
    from openpyxl import Workbook
    from docx import Document
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    # xlsx: bad (error value + placeholder) and good
    wb = Workbook(); ws = wb.active; ws["A1"] = "TBD"; ws["A2"] = "#DIV/0!"
    wb.save(tmp / "bad.xlsx")
    wb2 = Workbook(); wb2.active["A1"] = "clean"; wb2.save(tmp / "good.xlsx")
    (tmp / "data.csv").write_text("a,b\n1,2\n", encoding="utf-8")

    # xlsx with an UNCACHED error formula: openpyxl writes the formula but computes
    # nothing, so #DIV/0! is only detectable AFTER a real recalculation pass.
    wbf = Workbook(); wsf = wbf.active
    wsf["A1"] = 10; wsf["A2"] = 0; wsf["A3"] = "=A1/A2"
    wbf.save(tmp / "formula.xlsx")

    # docx: bad (tag) and good (heading)
    b = Document(); b.add_paragraph("Hi {{ name }}"); b.save(tmp / "bad.docx")
    g = Document(); g.add_heading("Title", 1); g.add_paragraph("ok"); g.save(tmp / "good.docx")

    # docx to templatize: a Label:value line + a bullet (for the template engine).
    t = Document(); t.add_heading("Weekly Update", 1)
    lp = t.add_paragraph(); lp.add_run("Client: ").bold = True; lp.add_run("Globex")
    t.add_paragraph("First achievement", style="List Bullet")
    t.save(tmp / "tmpl_src.docx")

    # pptx: bad (lorem) and good
    pb = Presentation(); s = pb.slides.add_slide(pb.slide_layouts[0])
    s.shapes.title.text = "Lorem ipsum"; pb.save(tmp / "bad.pptx")
    pg = Presentation(); sg = pg.slides.add_slide(pg.slide_layouts[0])
    sg.shapes.title.text = "Real Title"; pg.save(tmp / "good.pptx")

    # pptx to templatize: a title + TWO independent bullet lists (body placeholder +
    # a separate text box), to prove multi-list expansion doesn't interleave.
    from pptx.util import Inches
    pd = Presentation(); sd = pd.slides.add_slide(pd.slide_layouts[1])
    sd.shapes.title.text = "Globex"
    body = sd.placeholders[1].text_frame
    body.text = "Alpha one"; body.add_paragraph().text = "Alpha two"
    box = sd.shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(2)).text_frame
    box.text = "Beta one"; box.add_paragraph().text = "Beta two"
    pd.save(tmp / "tmpl_deck.pptx")

    # pdf: text layer and 'scanned' (rectangle only)
    c = canvas.Canvas(str(tmp / "text.pdf"), pagesize=letter)
    c.drawString(72, 720, "Real text content here."); c.showPage(); c.save()
    c2 = canvas.Canvas(str(tmp / "scan.pdf"), pagesize=letter)
    c2.rect(72, 600, 200, 100, fill=1); c2.showPage(); c2.save()

    # JSON output + schema for the structured-output validator (stdlib only)
    (tmp / "schema.json").write_text(
        json.dumps({"type": "object", "required": ["name"],
                    "properties": {"name": {"type": "string"}}}), encoding="utf-8")
    (tmp / "out_good.json").write_text('{"name":"x"}', encoding="utf-8")
    (tmp / "out_bad.json").write_text('{"nope":1}', encoding="utf-8")


def main():
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        make_fixtures(tmp)

        print("engineering-excel-workbooks/validate_workbook.py")
        rc, d = run("office/engineering-excel-workbooks/scripts/validate_workbook.py", tmp / "bad.xlsx")
        check("excel validator fails on bad workbook", rc == 1 and d and d["status"] == "FAILED")
        rc, d = run("office/engineering-excel-workbooks/scripts/validate_workbook.py", tmp / "good.xlsx")
        check("excel validator passes clean workbook", rc == 0 and d and d["status"] == "OK")

        print("engineering-excel-workbooks/recalculate_workbook.py")
        recalc_out = tmp / "formula.recalc.xlsx"
        rc, d = run("office/engineering-excel-workbooks/scripts/recalculate_workbook.py",
                    tmp / "formula.xlsx", recalc_out)
        status = (d or {}).get("status")
        if status == "RECALCULATED":
            # Real path (LibreOffice present): assert only the script CONTRACT (ran,
            # produced output). Whether the recompute populates cached error values is
            # environment-dependent (LibreOffice recalc-on-load policy) and is verified
            # MANUALLY, not asserted here — so CI stays green on the unverified branch
            # without over-claiming the gap is closed.
            name = "recalc runs and produces output (LibreOffice present)"
            ok = rc == 0 and bool(d.get("recalculated")) and recalc_out.exists()
        elif status == "SKIPPED_NO_ENGINE":
            # Fallback path (no LibreOffice here): must skip gracefully, NOT claim success.
            name = "recalc degrades gracefully when LibreOffice is absent"
            ok = rc == 0 and d.get("recalculated") is False
        else:
            name = f"recalc returned unexpected status: {status!r}"
            ok = False
        check(name, ok)

        print("processing-excel-files/extract_workbook.py")
        rc, d = run("office/processing-excel-files/scripts/extract_workbook.py", tmp / "data.csv")
        check("excel extractor reads csv", rc == 0 and d and d["format"] == "csv")

        print("processing-word-documents/extract_docx.py")
        rc, d = run("office/processing-word-documents/scripts/extract_docx.py", tmp / "good.docx")
        check("docx extractor returns markdown", rc == 0 and d and "Title" in d["markdown"])

        print("authoring-word-documents/validate_docx.py")
        rc, d = run("office/authoring-word-documents/scripts/validate_docx.py", tmp / "bad.docx")
        check("docx validator fails on unfilled tag", rc == 1 and d and d["status"] == "FAILED")
        rc, d = run("office/authoring-word-documents/scripts/validate_docx.py", tmp / "good.docx")
        check("docx validator passes clean doc", rc == 0 and d and d["status"] == "OK")

        print("processing-powerpoint-files/extract_pptx.py")
        rc, d = run("office/processing-powerpoint-files/scripts/extract_pptx.py", tmp / "good.pptx")
        check("pptx extractor returns slides", rc == 0 and d and d["slides"][0]["title"] == "Real Title")

        print("building-powerpoint-decks/validate_pptx.py")
        rc, d = run("office/building-powerpoint-decks/scripts/validate_pptx.py", tmp / "bad.pptx")
        check("pptx validator fails on lorem", rc == 1 and d and d["status"] == "FAILED")
        rc, d = run("office/building-powerpoint-decks/scripts/validate_pptx.py", tmp / "good.pptx")
        check("pptx validator passes clean deck", rc == 0 and d and d["status"] == "OK")

        print("processing-pdf-documents/extract_pdf.py")
        rc, d = run("office/processing-pdf-documents/scripts/extract_pdf.py", tmp / "text.pdf")
        check("pdf extractor finds text layer", rc == 0 and d and d["fidelity"]["has_text_layer"])
        rc, d = run("office/processing-pdf-documents/scripts/extract_pdf.py", tmp / "scan.pdf")
        check("pdf extractor flags scanned -> OCR", rc == 0 and d and not d["fidelity"]["has_text_layer"])

        print("processing-documents/detect_type.py")
        rc, d = run("office/processing-documents/scripts/detect_type.py", tmp / "good.docx")
        check("type detector identifies docx", rc == 0 and d and d["detected_type"] == "docx")

        print("building-document-templates: templatize -> build -> fill -> validate")
        tdir = "office/building-document-templates/scripts"
        reg = tmp / "registry"
        env = {"TEMPLATE_REGISTRY": reg}
        # 1. propose: extracts candidate fields (must find the labelled client value)
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_src.docx",
                    "--out", tmp / "proposal.json")
        prop = json.loads((tmp / "proposal.json").read_text(encoding="utf-8")) if rc == 0 else {}
        names = {c["current_text"] for c in prop.get("candidates", [])}
        check("templatize propose extracts candidates", rc == 0 and "Globex" in names)
        # 2. reviewed proposal: client_name (text) + achievements (list)
        (tmp / "reviewed.json").write_text(json.dumps({
            "format": "docx", "source_file": "tmpl_src.docx", "candidates": [
                {"current_text": "Globex", "suggest_name": "client_name",
                 "suggest_type": "text", "keep": "variable"},
                {"current_text": "First achievement", "suggest_name": "achievements",
                 "suggest_type": "list", "keep": "variable"},
            ]}), encoding="utf-8")
        # 3. build: registers template + manifest in the isolated registry
        rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_src.docx",
                    "--fields", tmp / "reviewed.json", "--client", "acme",
                    "--doc-type", "weekly", "--created", "2026-07-13", env=env)
        man = reg / "acme" / "weekly" / "manifest.json"
        check("templatize build registers template + manifest", rc == 0 and man.exists())
        # 4. fill: list field expands into real bullets
        (tmp / "content.json").write_text(json.dumps({
            "client_name": "Initech", "achievements": ["Alpha", "Beta", "Gamma"]}),
            encoding="utf-8")
        out_docx = tmp / "filled.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "weekly",
                    "--data", tmp / "content.json", "--out", out_docx, env=env)
        check("fill produces output document", rc == 0 and out_docx.exists())
        # 5. validate: filled doc OK, and a doc with a leftover tag FAILs
        rc, d = run(f"{tdir}/validate.py", out_docx)
        check("validate passes a fully filled document", rc == 0 and d and d["status"] == "OK")
        rc, d = run(f"{tdir}/validate.py", tmp / "bad.docx")
        check("validate fails on a leftover template tag", rc == 1 and d and d["status"] == "FAIL")
        # bonus: the list field became 3 bullet paragraphs
        from docx import Document as _Doc
        filled_texts = [p.text for p in _Doc(str(out_docx)).paragraphs]
        check("fill expands a list field into separate items",
              all(x in filled_texts for x in ("Alpha", "Beta", "Gamma")))

        # 6. missing REQUIRED field: tag stays, fill exits non-zero, validate FAILs.
        (tmp / "partial.json").write_text(json.dumps({"achievements": ["only this"]}),
                                          encoding="utf-8")  # client_name omitted
        out_partial = tmp / "partial.docx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "weekly",
                    "--data", tmp / "partial.json", "--out", out_partial, env=env)
        check("fill exits non-zero when a required field is missing",
              rc != 0 and out_partial.exists())
        rc, d = run(f"{tdir}/validate.py", out_partial)
        check("validate fails when a required field was left unfilled",
              rc == 1 and d and d["status"] == "FAIL")

        # PPTX round-trip with TWO independent list fields.
        rc, _ = run(f"{tdir}/templatize.py", "propose", "--file", tmp / "tmpl_deck.pptx",
                    "--out", tmp / "deck_prop.json")
        (tmp / "deck_reviewed.json").write_text(json.dumps({
            "format": "pptx", "source_file": "tmpl_deck.pptx", "candidates": [
                {"current_text": "Globex", "suggest_name": "client_name",
                 "suggest_type": "text", "keep": "variable"},
                {"current_text": "Alpha one", "suggest_name": "list_a",
                 "suggest_type": "list", "keep": "variable"},
                {"current_text": "Alpha two", "keep": "remove"},
                {"current_text": "Beta one", "suggest_name": "list_b",
                 "suggest_type": "list", "keep": "variable"},
                {"current_text": "Beta two", "keep": "remove"},
            ]}), encoding="utf-8")
        rc, _ = run(f"{tdir}/templatize.py", "build", "--file", tmp / "tmpl_deck.pptx",
                    "--fields", tmp / "deck_reviewed.json", "--client", "acme",
                    "--doc-type", "deck", "--created", "2026-07-13", env=env)
        deck_man = reg / "acme" / "deck" / "manifest.json"
        check("pptx templatize build registers template + manifest",
              rc == 0 and deck_man.exists())
        (tmp / "deck_content.json").write_text(json.dumps({
            "client_name": "Initech", "list_a": ["A1", "A2", "A3"],
            "list_b": ["B1", "B2"]}), encoding="utf-8")
        out_pptx = tmp / "filled.pptx"
        rc, _ = run(f"{tdir}/fill.py", "--client", "acme", "--doc-type", "deck",
                    "--data", tmp / "deck_content.json", "--out", out_pptx, env=env)
        check("pptx fill produces output deck", rc == 0 and out_pptx.exists())
        rc, d = run(f"{tdir}/validate.py", out_pptx)
        check("pptx validate passes a fully filled deck",
              rc == 0 and d and d["status"] == "OK")
        from pptx import Presentation as _Prs
        deck_texts = [C_para for sl in _Prs(str(out_pptx)).slides for sh in sl.shapes
                      if sh.has_text_frame for C_para in
                      [p.text for p in sh.text_frame.paragraphs]]
        check("pptx fill expands two independent list fields",
              all(x in deck_texts for x in ("A1", "A2", "A3", "B1", "B2")))

        print("generating-structured-outputs/validate_json_output.py")
        rc, d = run("ai-engineering/generating-structured-outputs/scripts/validate_json_output.py",
                    tmp / "out_good.json", tmp / "schema.json")
        check("json validator passes valid output", rc == 0 and d and d["status"] == "OK")
        rc, d = run("ai-engineering/generating-structured-outputs/scripts/validate_json_output.py",
                    tmp / "out_bad.json", tmp / "schema.json")
        check("json validator fails schema violation", rc == 1 and d and d["status"] == "FAILED")

    passed = sum(1 for _, ok in results if ok)
    total = len(results)
    print(f"\n{passed}/{total} checks passed")
    sys.exit(0 if passed == total else 1)


if __name__ == "__main__":
    main()
